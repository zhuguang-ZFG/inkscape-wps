"""导出为 Office/WPS 可打开的格式（docx/xlsx/pptx）及 Markdown（.md）。

策略（A+B 组合）：
- B：纯 Python 直写 docx/xlsx/pptx（默认可用，不依赖 soffice）。
- A：若检测到 LibreOffice（soffice），可在未来补充更高保真转换链（预留接口）。
"""

from __future__ import annotations

import csv
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional


class OfficeExportError(RuntimeError):
    pass


def _soffice_path() -> Optional[str]:
    return shutil.which("soffice") or shutil.which("libreoffice")


def has_soffice() -> bool:
    return _soffice_path() is not None


def _convert_with_soffice(input_path: Path, target_ext: str, output_path: Path) -> None:
    soff = _soffice_path()
    if not soff:
        raise OfficeExportError("未检测到 soffice。")
    with tempfile.TemporaryDirectory(prefix="inkscape-wps-export-") as td:
        outdir = Path(td)
        cmd = [
            soff,
            "--headless",
            "--convert-to",
            target_ext,
            "--outdir",
            str(outdir),
            str(input_path),
        ]
        try:
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        except subprocess.CalledProcessError as e:
            raise OfficeExportError(
                f"soffice 转换失败（{input_path.suffix} -> {target_ext}）：\n"
                f"{(e.stderr or b'').decode('utf-8', errors='replace')}"
            ) from e
        produced = outdir / (input_path.stem + "." + target_ext)
        if not produced.is_file():
            # 有些版本会改名，兜底搜同后缀
            candidates = list(outdir.glob(f"*.{target_ext}"))
            if not candidates:
                raise OfficeExportError(f"soffice 未生成目标文件：*.{target_ext}")
            produced = candidates[0]
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_bytes(produced.read_bytes())


@dataclass(frozen=True)
class DocRun:
    text: str
    bold: bool = False
    italic: bool = False
    underline: bool = False
    font_family: str | None = None
    font_pt: float | None = None


@dataclass(frozen=True)
class DocParagraph:
    runs: List[DocRun]
    align: str = "left"  # left/center/right/justify


def export_docx(
    path: Path,
    *,
    paragraphs: List[DocParagraph],
    html_text: str | None = None,
    prefer_soffice: bool = False,
) -> None:
    # A 路：优先用 soffice 把 HTML 转 DOCX（通常对段落版式更友好）
    if prefer_soffice and has_soffice() and html_text is not None:
        with tempfile.TemporaryDirectory(prefix="inkscape-wps-docx-") as td:
            src = Path(td) / "input.html"
            src.write_text(html_text, encoding="utf-8")
            _convert_with_soffice(src, "docx", path)
            return

    try:
        from docx import Document  # type: ignore
        from docx.enum.text import WD_ALIGN_PARAGRAPH  # type: ignore
    except Exception as e:
        raise OfficeExportError("缺少依赖：请安装 `python-docx` 后再导出 DOCX。") from e

    doc = Document()
    al_map = {
        "left": WD_ALIGN_PARAGRAPH.LEFT,
        "center": WD_ALIGN_PARAGRAPH.CENTER,
        "right": WD_ALIGN_PARAGRAPH.RIGHT,
        "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
    }
    for p in paragraphs:
        para = doc.add_paragraph()
        para.alignment = al_map.get(p.align, WD_ALIGN_PARAGRAPH.LEFT)
        for r in p.runs:
            run = para.add_run(r.text or "")
            run.bold = bool(r.bold)
            run.italic = bool(r.italic)
            run.underline = bool(r.underline)
            if r.font_family:
                run.font.name = r.font_family
            if r.font_pt and r.font_pt > 0:
                try:
                    from docx.shared import Pt  # type: ignore

                    run.font.size = Pt(float(r.font_pt))
                except Exception:
                    pass
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))


def export_xlsx(path: Path, *, table_blob: Dict[str, Any], prefer_soffice: bool = False) -> None:
    try:
        import openpyxl  # type: ignore
    except Exception as e:
        # A 路：若有 soffice，可先写 CSV 再转 XLSX（对 WPS 兼容较稳）
        if prefer_soffice and has_soffice():
            with tempfile.TemporaryDirectory(prefix="inkscape-wps-xlsx-") as td:
                src = Path(td) / "input.csv"
                cells = table_blob.get("cells") or []
                rows = int(table_blob.get("rows", len(cells) or 1))
                cols = int(table_blob.get("cols", (len(cells[0]) if cells else 1)))
                rows = max(1, rows)
                cols = max(1, cols)
                with src.open("w", encoding="utf-8", newline="") as f:
                    w = csv.writer(f)
                    for r in range(rows):
                        row_data = cells[r] if r < len(cells) and isinstance(cells[r], list) else []
                        line: list[str] = []
                        for c in range(cols):
                            cell = (
                                row_data[c]
                                if c < len(row_data) and isinstance(row_data[c], dict)
                                else {}
                            )
                            line.append(str(cell.get("text", "") or ""))
                        w.writerow(line)
                _convert_with_soffice(src, "xlsx", path)
                return
        raise OfficeExportError("缺少依赖：请安装 `openpyxl` 后再导出 XLSX。") from e

    def _mm_to_excel_col_width(mm: float) -> float:
        # 经验公式：像素≈ (宽度*7)+5；像素(mm) = mm*96/25.4
        px = float(mm) * 96.0 / 25.4
        return max(0.0, (px - 5.0) / 7.0)

    def _mm_to_excel_row_height_points(mm: float) -> float:
        # Excel row height 的单位为 points
        # 1 inch = 72 points = 25.4 mm
        return float(mm) * 72.0 / 25.4

    wb = openpyxl.Workbook()
    ws = wb.active
    cells = table_blob.get("cells") or []
    rows = int(table_blob.get("rows", len(cells) or 1))
    cols = int(table_blob.get("cols", (len(cells[0]) if cells else 1)))
    rows = max(1, rows)
    cols = max(1, cols)

    cell_w_mm = float(table_blob.get("cell_w_mm", 28.0))
    cell_h_mm = float(table_blob.get("cell_h_mm", 12.0))
    try:
        from openpyxl.utils import get_column_letter  # type: ignore
    except Exception:
        get_column_letter = None

    if get_column_letter is not None:
        col_w = _mm_to_excel_col_width(cell_w_mm)
        for c in range(cols):
            ws.column_dimensions[get_column_letter(c + 1)].width = float(col_w)

    row_h = _mm_to_excel_row_height_points(cell_h_mm)
    for r in range(rows):
        ws.row_dimensions[r + 1].height = float(row_h)

    for r in range(rows):
        row_data = cells[r] if r < len(cells) and isinstance(cells[r], list) else []
        for c in range(cols):
            cell = row_data[c] if c < len(row_data) and isinstance(row_data[c], dict) else {}
            ws.cell(row=r + 1, column=c + 1).value = str(cell.get("text", "") or "")
    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(path))


def export_pptx(path: Path, *, slides: List[str], prefer_soffice: bool = False) -> None:
    # 目前 PPTX 仍以 B 路为主；A 路可后续补 ODP 中转增强
    del prefer_soffice
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
    except Exception as e:
        raise OfficeExportError("缺少依赖：请安装 `python-pptx` 后再导出 PPTX。") from e

    prs = Presentation()
    # 使用 Title and Content 布局，保证 WPS 打开兼容性
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    for s in (slides or [""]):
        slide = prs.slides.add_slide(layout)
        # 标题留空，内容放到 body
        if slide.shapes.title is not None:
            slide.shapes.title.text = ""
        body = None
        for shp in slide.shapes:
            if shp.has_text_frame and shp is not slide.shapes.title:
                body = shp
                break
        if body is None:
            body = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        tf = body.text_frame
        tf.clear()
        lines = [ln for ln in str(s or "").splitlines() if ln.strip()]
        if not lines:
            lines = [""]
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ln
            try:
                p.font.size = Pt(20)
            except Exception:
                pass
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


def export_markdown(path: Path, *, body: str) -> None:
    """写入 UTF-8 Markdown 文本（无第三方依赖）。"""
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(body if body else "", encoding="utf-8")
