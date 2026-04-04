"""Office/WPS 文件导入（本机使用，尽量不引入硬依赖）。

支持：
- .docx → word_html（尽量保留基础段落/粗斜）
- .xlsx → table_blob（纯文本单元格）
- .pptx → slides（每页合并文本）
- .md / .markdown → 解析为纯文本（供单线编辑区；需 `markdown` 包）

WPS 私有格式：
- .wps/.et/.dps：尝试用 LibreOffice(soffice) 头less转换为 docx/xlsx/pptx 后再导入；
  若系统不存在 soffice，则抛出可读错误提示用户安装 LibreOffice 或用 WPS 另存为。
"""

from __future__ import annotations

import re
import shutil
import subprocess
import tempfile
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Dict, List, Optional


class OfficeImportError(RuntimeError):
    pass


def _require_import_file(path: Path) -> Path:
    """解析为绝对路径并确认存在且为常规文件，避免静默 OSError。"""
    p = path.expanduser()
    try:
        p = p.resolve()
    except OSError as e:
        raise OfficeImportError(f"无法解析路径：{path}") from e
    if not p.is_file():
        raise OfficeImportError(f"文件不存在或不是可读文件：{path}")
    return p


def _escape_html(s: str) -> str:
    return (
        s.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&#39;")
    )


def _runs_to_html(runs: list[tuple[str, bool, bool]]) -> str:
    out: list[str] = []
    for text, bold, italic in runs:
        t = _escape_html(text).replace("\n", "<br/>")
        if not t:
            continue
        if bold:
            t = f"<b>{t}</b>"
        if italic:
            t = f"<i>{t}</i>"
        out.append(t)
    return "".join(out)


def import_docx_to_html(path: Path) -> str:
    try:
        import docx  # type: ignore
    except ImportError as e:
        raise OfficeImportError(
            "缺少依赖：无法导入 .docx。请安装 `python-docx` 后重试。"
        ) from e

    fp = _require_import_file(Path(path))
    doc = docx.Document(str(fp))
    paras_html: list[str] = []
    for p in doc.paragraphs:
        runs: list[tuple[str, bool, bool]] = []
        for r in p.runs:
            txt = r.text or ""
            if not txt:
                continue
            runs.append((txt, bool(r.bold), bool(r.italic)))
        inner = _runs_to_html(runs)
        if not inner:
            paras_html.append('<p style="margin-top:0;margin-bottom:0;"><br/></p>')
        else:
            paras_html.append(f'<p style="margin-top:0;margin-bottom:0;">{inner}</p>')
    if not paras_html:
        return ""
    return "<html><body>" + "\n".join(paras_html) + "</body></html>"


def import_xlsx_to_table_blob(
    path: Path,
    *,
    max_rows: int = 200,
    max_cols: int = 50,
) -> Dict[str, Any]:
    try:
        import openpyxl  # type: ignore
    except ImportError as e:
        raise OfficeImportError(
            "缺少依赖：无法导入 .xlsx。请安装 `openpyxl` 后重试。"
        ) from e

    fp = _require_import_file(Path(path))
    # 注意：read_only=True 下 openpyxl 的 column_dimensions/row_dimensions 不可用，
    # 我们需要读取列宽/行高以反推 cell_w_mm/cell_h_mm，因此这里必须使用 read_only=False。
    wb = openpyxl.load_workbook(str(fp), read_only=False, data_only=True)
    ws = wb.active
    # 限制最大范围，避免超大表拖垮 UI
    max_r = min(int(ws.max_row or 1), int(max_rows))
    max_c = min(int(ws.max_column or 1), int(max_cols))

    def _excel_col_width_to_mm(width: float) -> float:
        # 经验公式：像素≈ (宽度*7)+5；像素(mm) = mm*96/25.4
        px = float(width) * 7.0 + 5.0
        return px * 25.4 / 96.0

    def _excel_row_height_to_mm(height_points: float) -> float:
        # points -> mm
        return float(height_points) * 25.4 / 72.0

    sheet_default_col_w = float(
        getattr(getattr(ws, "sheet_format", None), "defaultColWidth", 8.43)
        or 8.43
    )
    sheet_default_row_h_pt = float(
        getattr(getattr(ws, "sheet_format", None), "defaultRowHeight", 15.0)
        or 15.0
    )

    # 取行/列尺寸的平均值，作为统一 cell_w_mm/cell_h_mm。
    # 当前表格实现还不支持每列/每行独立尺寸。
    col_mm: list[float] = []
    for c in range(1, max_c + 1):
        letter = openpyxl.utils.get_column_letter(c)  # type: ignore[attr-defined]
        dim = ws.column_dimensions.get(letter)
        w = dim.width if dim is not None else None  # type: ignore[union-attr]
        if w is None:
            w = sheet_default_col_w
        try:
            mm = _excel_col_width_to_mm(float(w))
        except Exception:
            mm = _excel_col_width_to_mm(sheet_default_col_w)
        if mm > 0:
            col_mm.append(mm)
    row_mm: list[float] = []
    for r in range(1, max_r + 1):
        h = ws.row_dimensions.get(r).height if ws.row_dimensions.get(r) is not None else None  # type: ignore[union-attr]
        if h is None:
            h = sheet_default_row_h_pt
        try:
            mm = _excel_row_height_to_mm(float(h))
        except Exception:
            mm = _excel_row_height_to_mm(sheet_default_row_h_pt)
        if mm > 0:
            row_mm.append(mm)

    cell_w_mm = sum(col_mm) / len(col_mm) if col_mm else 28.0
    cell_h_mm = sum(row_mm) / len(row_mm) if row_mm else 12.0

    cells: list[list[dict[str, Any]]] = []
    for r in range(1, max_r + 1):
        row: list[dict[str, Any]] = []
        for c in range(1, max_c + 1):
            v = ws.cell(row=r, column=c).value
            s = "" if v is None else str(v)
            row.append({"text": s, "html": None})
        cells.append(row)
    return {
        "cell_w_mm": float(cell_w_mm),
        "cell_h_mm": float(cell_h_mm),
        "rows": max_r,
        "cols": max_c,
        "cells": cells,
    }


def import_pptx_to_slides(path: Path, *, max_slides: int = 200) -> List[str]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise OfficeImportError(
            "缺少依赖：无法导入 .pptx。请安装 `python-pptx` 后重试。"
        ) from e

    fp = _require_import_file(Path(path))
    prs = Presentation(str(fp))
    slides: list[str] = []
    for i, s in enumerate(prs.slides):
        if i >= max_slides:
            break
        parts: list[str] = []
        for shp in s.shapes:
            if not hasattr(shp, "text"):
                continue
            txt = (shp.text or "").strip()
            if txt:
                parts.append(txt)
        slides.append("\n".join(parts).strip())
    return slides if slides else [""]


def _soffice_path() -> Optional[str]:
    return shutil.which("soffice") or shutil.which("libreoffice")


def try_convert_wps_private_to_office(path: Path) -> Path:
    """把 .wps/.et/.dps 转成 docx/xlsx/pptx，返回转换后的新文件路径。"""
    _require_import_file(Path(path))
    ext = path.suffix.lower()
    if ext not in (".wps", ".et", ".dps"):
        return path
    soff = _soffice_path()
    if not soff:
        raise OfficeImportError(
            "无法导入 WPS 私有格式（.wps/.et/.dps）：未找到 LibreOffice(soffice)。\n"
            "解决方案：安装 LibreOffice，或用 WPS 手动“另存为” docx/xlsx/pptx。"
        )
    # 映射到目标格式
    target_ext = {".wps": "docx", ".et": "xlsx", ".dps": "pptx"}[ext]
    with tempfile.TemporaryDirectory(prefix="inkscape-wps-convert-") as td:
        outdir = Path(td)
        cmd = [soff, "--headless", "--convert-to", target_ext, "--outdir", str(outdir), str(path)]
        try:
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        except subprocess.CalledProcessError as e:
            raise OfficeImportError(
                f"LibreOffice 转换失败（{ext} → {target_ext}）。\n"
                f"stderr:\n{(e.stderr or b'').decode('utf-8', errors='replace')}"
            ) from e
        # 找到转换产物
        outs = list(outdir.glob(f"*.{target_ext}"))
        if not outs:
            raise OfficeImportError("LibreOffice 未生成转换文件。")
        # 拷贝到与源文件同目录，便于后续“最近打开”
        dst = path.with_suffix("." + target_ext)
        try:
            dst.write_bytes(outs[0].read_bytes())
        except OSError as e:
            raise OfficeImportError(f"写入转换文件失败：{dst}") from e
        return dst


def detect_office_kind(path: Path) -> str:
    """返回：docx/xlsx/pptx/md/wps/et/dps/unknown"""
    ext = path.suffix.lower().lstrip(".")
    if ext in ("docx", "xlsx", "pptx", "wps", "et", "dps"):
        return ext
    if ext in ("md", "markdown"):
        return "md"
    return "unknown"


class _HtmlToPlainParser(HTMLParser):
    """将常见块级 HTML 转为带换行的纯文本（用于 Markdown 渲染结果）。"""

    _BLOCK = frozenset(
        {
            "p",
            "div",
            "br",
            "h1",
            "h2",
            "h3",
            "h4",
            "h5",
            "h6",
            "li",
            "tr",
            "blockquote",
            "pre",
            "table",
            "ul",
            "ol",
            "thead",
            "tbody",
            "hr",
        }
    )

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: Any) -> None:
        if tag == "br" or tag == "hr":
            self._parts.append("\n")
        elif tag in self._BLOCK:
            if self._parts and not str(self._parts[-1]).endswith("\n"):
                self._parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in self._BLOCK and tag not in ("br", "hr"):
            self._parts.append("\n")

    def handle_data(self, data: str) -> None:
        self._parts.append(data)

    def plain(self) -> str:
        text = "".join(self._parts)
        text = re.sub(r"[ \t]+\n", "\n", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()


def html_fragment_to_plain_text(html: str) -> str:
    """把 HTML 片段（如 Markdown 渲染结果）压成纯文本段落。"""
    p = _HtmlToPlainParser()
    p.feed(f"<div>{html}</div>")
    p.close()
    return p.plain()


def import_markdown_string_to_plain(text: str) -> str:
    """
    将 Markdown 源码转为纯文本（块级换行尽量保留）。
    依赖 `markdown` 包。
    """
    raw = text if text else ""
    try:
        import markdown as md_pkg  # type: ignore
    except ImportError as e:
        raise OfficeImportError(
            "缺少依赖：无法解析 Markdown。请安装 `markdown` 后重试（pip install markdown）。"
        ) from e

    extensions = ["extra", "nl2br"]
    try:
        html = md_pkg.markdown(raw, extensions=extensions)
    except Exception as e:
        raise OfficeImportError(f"Markdown 解析失败：{e}") from e
    return html_fragment_to_plain_text(html) if html.strip() else ""


def import_markdown_to_plain(path: Path) -> str:
    """读取 .md 文件并 `import_markdown_string_to_plain`。"""
    fp = _require_import_file(Path(path))
    return import_markdown_string_to_plain(fp.read_text(encoding="utf-8", errors="replace"))


_MD_SLIDE_SPLIT = re.compile(r"^\s*---\s*$", re.MULTILINE)


def _strip_yaml_front_matter(text: str) -> str:
    """去掉开头的 YAML front matter（--- … ---），避免误当作分页。"""
    t = text.lstrip("\ufeff").lstrip()
    if not t.startswith("---"):
        return text
    rest = t[3:].lstrip("\n")
    m = _MD_SLIDE_SPLIT.search(rest)
    if not m:
        return text
    return rest[m.end() :].lstrip()


def split_markdown_into_slides(raw: str) -> Optional[List[str]]:
    """
    若正文含独立一行的 --- 分隔符，则拆成多段；否则返回 None。
    用于与「演示」页分页习惯对齐（Markdown 常用 --- 作页分隔）。
    """
    s = _strip_yaml_front_matter(raw or "").strip()
    if not s:
        return None
    parts = [p.strip() for p in _MD_SLIDE_SPLIT.split(s) if p.strip()]
    if len(parts) < 2:
        return None
    return parts


def import_markdown_file_to_slides_plain(path: Path) -> Optional[List[str]]:
    """若文件含 --- 分页则返回每页纯文本列表，否则 None。"""
    fp = _require_import_file(Path(path))
    raw = fp.read_text(encoding="utf-8", errors="replace")
    chunks = split_markdown_into_slides(raw)
    if chunks is None:
        return None
    return [import_markdown_string_to_plain(c) for c in chunks]
